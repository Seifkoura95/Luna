"""
Build training assets for the Luna Group VIP Staff Training SOP.

Outputs:
  - /app/LUNA_STAFF_TRAINING_SOP.pdf       (cover page + full SOP)
  - /app/LUNA_STAFF_TRAINING_DECK.pptx     (~2-hour facilitator deck)
"""

import os
import re
import markdown
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from weasyprint import HTML, CSS

ROOT = "/app"
SOP_MD = f"{ROOT}/LUNA_STAFF_TRAINING_SOP.md"
PDF_OUT = f"{ROOT}/LUNA_STAFF_TRAINING_SOP.pdf"
PPTX_OUT = f"{ROOT}/LUNA_STAFF_TRAINING_DECK.pptx"
LOGO_URL = "https://customer-assets.emergentagent.com/job_c826baa4-6640-40ce-9e0d-38132d9944fc/artifacts/2k76js5m_luna-group-logo-2.webp"

# =========================================================================
# 1) PDF — cover page + rendered Markdown
# =========================================================================

with open(SOP_MD, "r") as f:
    md_text = f.read()

# Strip the H1 from the SOP body (we put a custom cover instead)
md_body = re.sub(r"^#\s+LUNA GROUP VIP — STAFF TRAINING SOP\s*\n", "", md_text, count=1)

html_body = markdown.markdown(
    md_body,
    extensions=["tables", "fenced_code", "toc", "sane_lists"],
)

cover_html = f"""
<div class="cover">
  <img src="{LOGO_URL}" class="logo" />
  <div class="cover-tagline">BRISBANE • GOLD COAST</div>
  <div class="cover-title">STAFF TRAINING<br/>STANDARD OPERATING PROCEDURE</div>
  <div class="cover-sub">Luna Group VIP App · v1.0 · February 2026</div>
  <div class="cover-foot">
    <div>CONFIDENTIAL — INTERNAL USE ONLY</div>
    <div>Door · Bar · Floor · Supervisor · Manager · Admin</div>
  </div>
</div>
<div style="page-break-after: always;"></div>
"""

full_html = f"""<!doctype html>
<html><head><meta charset="utf-8"><title>Luna Staff Training SOP</title></head>
<body>
{cover_html}
<main>{html_body}</main>
</body></html>
"""

pdf_css = CSS(string=r"""
@page {
    size: A4;
    margin: 22mm 18mm 22mm 18mm;
    @bottom-center {
        content: "Luna Group VIP — Staff Training SOP   ·   Page " counter(page) " of " counter(pages);
        font-family: 'Helvetica Neue', Arial, sans-serif;
        font-size: 8pt;
        color: #666;
    }
}
@page :first { margin: 0; @bottom-center { content: ""; } }

body {
    font-family: 'Helvetica Neue', Arial, sans-serif;
    color: #1a1a1f;
    line-height: 1.55;
    font-size: 10.5pt;
}

.cover {
    background: linear-gradient(155deg, #08080E 0%, #1a1530 60%, #0a0a14 100%);
    color: #f0f0f8;
    height: 297mm;
    padding: 60mm 25mm 30mm 25mm;
    box-sizing: border-box;
    position: relative;
}
.cover .logo { width: 110mm; display: block; margin: 0 auto 18mm auto; }
.cover-tagline {
    text-align: center; letter-spacing: 6px; font-size: 11pt;
    color: #D4AF5A; margin-bottom: 35mm;
}
.cover-title {
    text-align: center; font-size: 30pt; font-weight: 800;
    letter-spacing: 2px; line-height: 1.2;
    color: #ffffff; margin-bottom: 8mm;
}
.cover-sub {
    text-align: center; font-size: 12pt; letter-spacing: 2px;
    color: #D4AF5A; opacity: 0.85;
}
.cover-foot {
    position: absolute; left: 25mm; right: 25mm; bottom: 25mm;
    display: flex; justify-content: space-between;
    font-size: 8.5pt; letter-spacing: 1.5px;
    color: rgba(240,240,248,0.55);
    border-top: 1px solid rgba(212,175,90,0.4);
    padding-top: 8mm;
}

main { padding-top: 0; }

h1 {
    font-size: 20pt; color: #1a1a1f;
    border-bottom: 2px solid #D4AF5A;
    padding-bottom: 4mm;
    margin-top: 8mm;
}
h2 {
    font-size: 14pt; color: #08080E;
    margin-top: 9mm; margin-bottom: 3mm;
    page-break-after: avoid;
}
h3 { font-size: 11.5pt; color: #2a2540; margin-top: 6mm; }
h4 { font-size: 10.5pt; color: #4a4360; }

p { margin: 0 0 3mm 0; }

table {
    width: 100%; border-collapse: collapse; margin: 3mm 0 5mm 0;
    font-size: 9.5pt;
    page-break-inside: avoid;
}
th {
    background: #1a1530; color: #D4AF5A;
    text-align: left; padding: 2mm 3mm;
    font-weight: 600; letter-spacing: 0.5px;
    border: 1px solid #D4AF5A;
}
td { padding: 2mm 3mm; border: 1px solid #d8d4e6; vertical-align: top; }
tr:nth-child(even) td { background: #faf9fc; }

ul, ol { padding-left: 7mm; margin: 0 0 4mm 0; }
li { margin: 0.8mm 0; }

code {
    background: #f0eef7;
    padding: 0.5mm 1.5mm; border-radius: 1mm;
    font-size: 9pt; color: #2a2540;
}

blockquote {
    border-left: 3px solid #D4AF5A;
    background: #fdfaf2;
    margin: 3mm 0; padding: 3mm 5mm;
    font-style: italic; color: #4a4360;
}

hr { border: none; border-top: 1px solid #d8d4e6; margin: 6mm 0; }

a { color: #5b3aa0; text-decoration: none; }
""")

print("Building PDF…")
HTML(string=full_html, base_url=ROOT).write_pdf(PDF_OUT, stylesheets=[pdf_css])
print(f"  → {PDF_OUT} ({os.path.getsize(PDF_OUT)//1024} KB)")

# =========================================================================
# 2) PPTX — 2-hour facilitator deck
# =========================================================================
print("Building PPTX deck…")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY      = RGBColor(0x08, 0x08, 0x0E)
DEEP      = RGBColor(0x1A, 0x15, 0x30)
GOLD      = RGBColor(0xD4, 0xAF, 0x5A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SOFT      = RGBColor(0xF0, 0xF0, 0xF8)
MUTED     = RGBColor(0x9A, 0x95, 0xB0)
RED       = RGBColor(0xE3, 0x18, 0x37)


def fill_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, *,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             font_name="Helvetica"):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font_name
    return tx


def add_bullets(slide, bullets, left, top, width, height, *,
                size=18, color=SOFT, bullet_color=GOLD, font_name="Helvetica"):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        # bullet marker
        rb = p.add_run()
        rb.text = "•   "
        rb.font.size = Pt(size)
        rb.font.bold = True
        rb.font.color.rgb = bullet_color
        rb.font.name = font_name
        # bullet text
        rt = p.add_run()
        rt.text = b
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
        rt.font.name = font_name
    return tx


def header_band(slide, kicker, title, *, kicker_color=GOLD):
    add_text(slide, kicker, 0.6, 0.5, 12, 0.4,
             size=12, bold=True, color=kicker_color)
    add_text(slide, title, 0.6, 0.95, 12, 1.0,
             size=32, bold=True, color=WHITE)
    # gold underline
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.85), Inches(1.0), Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def footer_band(slide, mod_label):
    add_text(slide, "Luna Group VIP · Staff Training", 0.6, 7.05, 6, 0.3,
             size=9, color=MUTED)
    add_text(slide, mod_label, 7.5, 7.05, 5.2, 0.3,
             size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    return slide


# ---------- Slide 1: Cover ----------
s = add_blank_slide()
fill_bg(s, NAVY)
# decorative gold dot
dot = s.shapes.add_shape(9, Inches(11.6), Inches(0.6), Inches(0.5), Inches(0.5))
dot.fill.solid(); dot.fill.fore_color.rgb = GOLD
dot.line.fill.background()

add_text(s, "LUNA GROUP", 0.7, 2.4, 12, 1.2,
         size=64, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "STAFF TRAINING", 0.7, 3.6, 12, 0.8,
         size=36, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
add_text(s, "Facilitator Deck · 2-hour module", 0.7, 4.4, 12, 0.5,
         size=18, color=SOFT, align=PP_ALIGN.LEFT)
add_text(s, "BRISBANE  ·  GOLD COAST", 0.7, 6.4, 12, 0.4,
         size=12, color=MUTED, align=PP_ALIGN.LEFT)
add_text(s, "v1.2 · Apr 2026", 0.7, 6.8, 12, 0.4,
         size=10, color=MUTED, align=PP_ALIGN.LEFT)

# ---------- Slide 2: Agenda ----------
s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "AGENDA", "What we'll cover today")
add_bullets(s, [
    "Module 1 · Welcome & The 9 Venues  (10 min)",
    "Module 2 · Roles & Access Matrix  (10 min)",
    "Module 3 · Member App Tour  (20 min)",
    "Module 4 · Staff Portal — Award · Validate · History  (30 min)",
    "Module 5 · QR Codes Decoded  (10 min)",
    "Module 6 · Points · Tiers · Subscriptions  (15 min)",
    "Module 7 · Safety / SOS Workflow  (10 min)",
    "Module 8 · Troubleshooting & Cheat Sheet  (10 min)",
    "Module 9 · Practical Drills + Sign-off  (15 min)",
], 0.7, 2.2, 12, 5, size=18)
footer_band(s, "Agenda")

# ---------- Slide 3: Module 1 divider ----------
def divider(num, title):
    s = add_blank_slide(); fill_bg(s, DEEP)
    add_text(s, f"MODULE {num}", 0.7, 2.6, 12, 0.6,
             size=18, bold=True, color=GOLD)
    add_text(s, title, 0.7, 3.2, 12, 1.5,
             size=44, bold=True, color=WHITE)
    line = s.shapes.add_shape(1, Inches(0.7), Inches(4.9), Inches(1.5), Inches(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    return s

divider(1, "Welcome &\nThe 9 Venues")

# Slide: 9 venues
s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 1", "The 9 Luna Group Venues")
venues = [
    ("Eclipse",       "Brisbane"),
    ("After Dark",    "Brisbane"),
    ("Pump",          "Brisbane"),
    ("Mamacita",      "Brisbane"),
    ("Juju",          "Brisbane"),
    ("Night Market",  "Brisbane"),
    ("Ember & Ash",   "Brisbane"),
    ("Su Casa BNE",   "Brisbane"),
    ("Su Casa GC",    "Gold Coast"),
]
# 3x3 grid
for i, (name, city) in enumerate(venues):
    col, row = i % 3, i // 3
    x = 0.7 + col * 4.15
    y = 2.2 + row * 1.55
    box = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(3.95), Inches(1.35))
    box.fill.solid(); box.fill.fore_color.rgb = DEEP
    box.line.color.rgb = GOLD; box.line.width = Pt(0.5)
    add_text(s, name, x + 0.2, y + 0.2, 3.6, 0.5, size=20, bold=True, color=WHITE)
    add_text(s, city.upper(), x + 0.2, y + 0.78, 3.6, 0.4, size=10, color=GOLD)
footer_band(s, "Module 1 · Venues")

# Slide: 3 pillars
s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 1", "The 3 Pillars of the Luna App")
pillars = [
    ("EARN",    "Every dollar a guest spends becomes points.\nSwiftPOS is the source of truth."),
    ("ENJOY",   "Drinks, entries, bottle service, VIP tables,\nevents, auctions."),
    ("ENGAGE",  "Luna AI concierge, leaderboards, missions,\nmilestones, crews."),
]
for i, (h, body) in enumerate(pillars):
    x = 0.7 + i * 4.15
    box = s.shapes.add_shape(5, Inches(x), Inches(2.4), Inches(3.95), Inches(3.6))
    box.fill.solid(); box.fill.fore_color.rgb = DEEP
    box.line.color.rgb = GOLD; box.line.width = Pt(0.5)
    add_text(s, h, x + 0.3, 2.7, 3.5, 0.7, size=28, bold=True, color=GOLD)
    add_text(s, body, x + 0.3, 3.6, 3.5, 2.3, size=14, color=SOFT)
add_text(s, "PRINCIPLE — never tell a guest \"the points didn't go through\". Always check Staff Portal first.",
         0.7, 6.3, 12, 0.6, size=12, color=GOLD, bold=True)
footer_band(s, "Module 1 · Pillars")

# ---------- Module 2 — Roles ----------
divider(2, "Roles &\nAccess Matrix")

s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 2", "Who can do what?")
roles_table = [
    ["Capability",                    "Member", "Staff", "Manager", "Admin"],
    ["Sign in to consumer app",       "✓", "✓", "✓", "✓"],
    ["Earn / spend points",           "✓", "✓", "✓", "✓"],
    ["Access Staff Portal",           "—", "✓", "✓", "✓"],
    ["Quick-Award points",            "—", "✓", "✓", "✓"],
    ["Validate reward / entry QRs",   "—", "✓", "✓", "✓"],
    ["Gift points (override)",        "—", "—", "✓", "✓"],
    ["Venue Dashboard",               "—", "—", "✓", "✓"],
    ["Manage menus / events",         "—", "—", "✓", "✓"],
    ["Reports (SwiftPOS / payments)", "—", "—", "limited", "✓"],
    ["Manage users / roles",          "—", "—", "—", "✓"],
]
rows, cols = len(roles_table), len(roles_table[0])
left, top, w, h = Inches(0.7), Inches(2.2), Inches(12), Inches(4.6)
table = s.shapes.add_table(rows, cols, left, top, w, h).table
col_widths = [Inches(5), Inches(1.75), Inches(1.75), Inches(1.75), Inches(1.75)]
for i, cw in enumerate(col_widths):
    table.columns[i].width = cw
for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = roles_table[r][c]
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(12)
                run.font.name = "Helvetica"
                run.font.color.rgb = WHITE if r == 0 else SOFT
                run.font.bold = (r == 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DEEP if r == 0 else NAVY
footer_band(s, "Module 2 · Roles")

# ---------- Module 3 — Member App ----------
divider(3, "Member App Tour")

s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 3", "5 bottom tabs you'll see")
tabs_info = [
    ("TONIGHT",  "Live events, AI picks,\nauctions, quick links"),
    ("VENUES",   "Map + list of all 9 venues,\ndetail screens, bookings"),
    ("WALLET",   "Points balance, history,\nMember Card QR, rewards"),
    ("LUNA AI",  "Personalised concierge,\nplan-a-night, recommendations"),
    ("PROFILE",  "Settings, missions,\nmilestones, subscriptions, safety"),
]
for i, (h, body) in enumerate(tabs_info):
    x = 0.7 + i * 2.5
    box = s.shapes.add_shape(5, Inches(x), Inches(2.3), Inches(2.35), Inches(3.8))
    box.fill.solid(); box.fill.fore_color.rgb = DEEP
    box.line.color.rgb = GOLD; box.line.width = Pt(0.5)
    add_text(s, h, x + 0.15, 2.5, 2.1, 0.6, size=16, bold=True, color=GOLD)
    add_text(s, body, x + 0.15, 3.2, 2.1, 2.6, size=12, color=SOFT)
footer_band(s, "Module 3 · App tour")

# Slide: detail screens
s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 3", "Detail screens worth knowing")
add_bullets(s, [
    "Member Card  →  full-screen QR for door + bar",
    "Rewards Shop  →  redeem points for vouchers / drinks / entries",
    "Bottle Service  →  packages bookable in-app",
    "Events  →  RSVP + buy tickets",
    "Auctions  →  bid on premium experiences",
    "Missions / Milestones  →  bonus points for activity",
    "Birthday Club  →  birthday-week perks",
    "Refer a Friend  →  shared bonus on first visit",
    "Lost & Found  →  report or claim",
    "Safety  →  one-tap SOS + emergency contacts",
], 0.7, 2.2, 12, 5, size=16)
footer_band(s, "Module 3 · Stack screens")

# ---------- Module 4 — Staff Portal ----------
divider(4, "Staff Portal\nAward · Validate · History")

# Slide: Award flow
s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 4 · TAB 1", "Award Points — only when SwiftPOS can't")
add_bullets(s, [
    "1. Confirm the correct venue is selected (top of screen)",
    "2. Find the member: search OR scan their Member Card",
    "3. Member profile loads (tier · balance · today)",
    "4. Enter dollar amount spent (max $50,000 / txn)",
    "5. Pick category: Drinks · Food · Entry · Booth · Bottles · Other",
    "6. ALWAYS enter the receipt reference — anti-double-credit guard",
    "7. Tap Award Points  →  $1 = 10 pts × tier multiplier",
    "8. 409 \"already awarded\"?  →  it's already in SwiftPOS — leave it",
], 0.7, 2.1, 12, 5, size=17)
footer_band(s, "Module 4 · Award")

# NEW slide — When to use vs not use Quick Award (architectural)
s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 4 · IMPORTANT", "Quick Award is an EXCEPTION tool")
add_text(s, "Use Quick Award only when SwiftPOS can't capture the sale.",
         0.7, 2.1, 12, 0.5, size=15, color=GOLD, bold=True)
add_bullets(s, [
    "DO  →  Cash-only / pop-up / off-system events",
    "DO  →  SwiftPOS outage / network down (note receipt, enter when back)",
    "DO  →  Split bills — only for the portion that didn't go through POS",
    "DO  →  Manager Gift Points — comps, hospitality, artist drinks",
    "",
    "DO NOT  →  Sale that already went through SwiftPOS (= double credit)",
    "DO NOT  →  Member 'forgot to give number' on a SwiftPOS docket",
    "             (manager attaches member to existing receipt instead)",
    "",
    "Receipt-ref guard returns 409 if the same receipt is awarded twice.",
], 0.7, 2.7, 12, 5, size=14)
footer_band(s, "Module 4 · Quick-Award rules")

# Slide: Validate flow
s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 4 · TAB 2", "Validate — door / bar / cashier")
add_bullets(s, [
    "Confirm venue selector",
    "Tap Scan QR  OR  paste the QR string into the input",
    "The portal auto-routes by prefix (next slide)",
    "GREEN = success — proceed; reward is now consumed",
    "RED = failure — read the message (expired / wrong venue / used)",
    "NEVER strip the LUNA- prefix",
    "If a refusal needs reversing AFTER green: manager only",
], 0.7, 2.2, 12, 5, size=18)
footer_band(s, "Module 4 · Validate")

# Slide: History
s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 4 · TAB 3", "History — your audit trail")
add_bullets(s, [
    "Filter: today · 7-day · 30-day · per-venue",
    "Lists every Award + every Validate by you and your team",
    "Use this FIRST to investigate any 'missing points' dispute",
    "Tap a row to see receipt reference, amount, member, timestamp",
    "Manager Gift-Points are flagged in red — bypass the earn-guard",
], 0.7, 2.2, 12, 5, size=18)
footer_band(s, "Module 4 · History")

# ---------- Module 5 — QR codes ----------
divider(5, "QR Codes Decoded")

s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 5", "5 prefixes to recognise instantly")
qr_table = [
    ["Prefix",          "Meaning",            "Where to scan"],
    ["LUNA-MEMBER:",    "Member ID card",     "AWARD tab"],
    ["LUNA-ENT-",       "Free Entry pass",    "VALIDATE — door"],
    ["LUNA-TKT-",       "Milestone ticket",   "VALIDATE — door"],
    ["LUNA-DRINK-",     "Free drink reward",  "VALIDATE — bar"],
    ["LUNA-BDAY-",      "Birthday reward",    "VALIDATE"],
    ["LUNA-<uuid>-…",   "Generic reward",     "VALIDATE"],
]
left, top, w, h = Inches(0.7), Inches(2.2), Inches(12), Inches(3.8)
t = s.shapes.add_table(len(qr_table), 3, left, top, w, h).table
t.columns[0].width = Inches(3.3)
t.columns[1].width = Inches(4.5)
t.columns[2].width = Inches(4.2)
for r, row in enumerate(qr_table):
    for c, val in enumerate(row):
        cell = t.cell(r, c)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.bold = (r == 0)
                run.font.color.rgb = WHITE if r == 0 else SOFT
                run.font.name = "Courier New" if (r > 0 and c == 0) else "Helvetica"
        cell.fill.solid()
        cell.fill.fore_color.rgb = DEEP if r == 0 else NAVY
add_text(s, "If it doesn't start with LUNA- it's not ours. Reject politely.",
         0.7, 6.3, 12, 0.6, size=14, color=GOLD, bold=True)
footer_band(s, "Module 5 · QR")

# ---------- Module 6 — Tiers & subs ----------
divider(6, "Points · Tiers · Subscriptions")

s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 6", "Base Tiers (lifetime points)")
tier_table = [
    ["Tier", "Lifetime points", "Multiplier", "Headline perks"],
    ["Bronze",  "0",       "1.0×", "Standard rewards"],
    ["Silver",  "5,000",   "1.25×","Priority entry occasionally"],
    ["Gold",    "25,000",  "1.5×", "Free drinks · select line-skips"],
    ["Legend",  "100,000", "2.0×", "All-access · dedicated host"],
]
left, top, w, h = Inches(0.7), Inches(2.3), Inches(12), Inches(3.5)
t = s.shapes.add_table(len(tier_table), 4, left, top, w, h).table
widths = [2, 3, 2.5, 4.5]
for i, ww in enumerate(widths):
    t.columns[i].width = Inches(ww)
for r, row in enumerate(tier_table):
    for c, val in enumerate(row):
        cell = t.cell(r, c)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(14)
                run.font.bold = (r == 0)
                run.font.color.rgb = WHITE if r == 0 else SOFT
                run.font.name = "Helvetica"
        cell.fill.solid()
        cell.fill.fore_color.rgb = DEEP if r == 0 else NAVY
footer_band(s, "Module 6 · Tiers")

s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 6", "Paid Subscriptions — Stripe")
add_bullets(s, [
    "Aurora  —  bonus points multiplier · early access",
    "Lunar   —  monthly free entries (LUNA-ENT-) · monthly drink · priority bookings",
    "Legend  —  Lunar perks + table credits + concierge",
    "",
    "Subscription tier is INDEPENDENT of base tier",
    "(A Bronze member can hold a Lunar subscription)",
    "",
    "Total earn multiplier = base tier × subscription bonus × active boost",
], 0.7, 2.2, 12, 5, size=18)
footer_band(s, "Module 6 · Subs")

# ── Missions & Milestones (April 2026 update) ────────────────────────
s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 6.5 · NEW", "Missions & Milestones — what changed")
add_bullets(s, [
    "MISSIONS auto-progress every time you Quick-Award. No extra steps.",
    "Members see live progress in Profile → Missions",
    "Mission complete? — push fires; member taps Claim to bank the points",
    "MILESTONES are lifetime tiers — Newbie · Rising Star · VIP Status · Luna Elite · Supernova · Legend",
    "Hitting a milestone unlocks signed LUNA-TKT- QRs (single-use, can't be forged)",
    "Manager-only: Lovable now has a 'Test trigger' button + activity timeline",
    "Common Q: \"did this count toward my mission?\" — Answer: yes, automatically",
], 0.7, 2.2, 12, 5, size=17)
footer_band(s, "Module 6.5 · Missions")

# ---------- Module 7 — Safety ----------
divider(7, "Safety / SOS\nWorkflow")

s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 7", "If a guest fires SOS in your venue")
add_bullets(s, [
    "Manager receives instant alert in the Lovable admin portal",
    "Locate guest — their last known venue is in the alert payload",
    "Provide assistance: water · first aid · ride-share · escort",
    "Write a 2-line incident note in the portal Resolve action",
    "If serious (medical / police): call 000 FIRST, then update portal",
    "",
    "Members add up to 5 emergency contacts in Safety Settings",
    "SMS dispatch via Twilio — coming next sprint",
], 0.7, 2.2, 12, 5, size=18)
footer_band(s, "Module 7 · Safety")

# ---------- Module 8 — Troubleshooting & cheat sheet ----------
divider(8, "Troubleshooting\n& Cheat Sheet")

s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 8", "Top member questions — one-line scripts")
qa = [
    ("\"How do I get my points?\"",          "I'll scan your member QR and award them now."),
    ("\"Where's my member QR?\"",            "Open the app → Wallet → Member Card."),
    ("\"I didn't get points from yesterday\"", "Let me check History — SwiftPOS can take a few minutes."),
    ("\"How do I get free entry?\"",         "Lunar subscription gives monthly entries; tier may also unlock."),
    ("\"What's my tier?\"",                  "It's on your Profile screen — Bronze · Silver · Gold · Legend."),
    ("\"How do referrals work?\"",           "Profile → Refer a Friend. You both get bonus points after their first visit."),
    ("\"I lost something\"",                 "App → Lost & Found → Report Missing. We'll match it."),
    ("\"Cancel my subscription?\"",          "Profile → Subscriptions → Manage. Stripe handles it; access until period-end."),
]
top = 2.2
for q, a in qa:
    add_text(s, q, 0.7, top, 5.5, 0.4, size=13, bold=True, color=GOLD)
    add_text(s, a, 6.3, top, 6.5, 0.4, size=13, color=SOFT)
    top += 0.55
footer_band(s, "Module 8 · Cheat sheet")

s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 8", "Troubleshooting playbook")
add_bullets(s, [
    "Points not showing  →  check History; if absent, re-award + flag manager",
    "QR won't scan  →  brightness up · full-screen Member Card · or search by email",
    "App crashing  →  force close · update · capture screenshot for support",
    "Wrong venue on receipt  →  do NOT reverse · call manager · they fix in Lovable",
    "Duplicate award  →  manager reverses one in Lovable Admin Portal",
], 0.7, 2.2, 12, 5, size=18)
footer_band(s, "Module 8 · Troubleshooting")

# ---------- Module 9 — drills + sign-off ----------
divider(9, "Practical Drills\n& Sign-off")

s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 9", "Drills — pair up and demonstrate")
add_bullets(s, [
    "DRILL 1 · Award $85 of drinks at Eclipse with receipt ref",
    "DRILL 2 · Validate a LUNA-DRINK- QR at the bar",
    "DRILL 3 · Validate a LUNA-ENT- QR at the door (wrong venue first → see error)",
    "DRILL 4 · Search a member by phone, then by QR scan",
    "DRILL 5 · Investigate a fake \"missing points\" complaint via History",
    "DRILL 6 · Manager only: issue a Gift-Points override with a written reason",
], 0.7, 2.2, 12, 5, size=18)
footer_band(s, "Module 9 · Drills")

s = add_blank_slide(); fill_bg(s, NAVY)
header_band(s, "MODULE 9", "Sign-off — each trainee must demonstrate")
add_bullets(s, [
    "✓ Log in and select correct venue",
    "✓ Search by name · email · QR scan (all three)",
    "✓ Quick-award with category + receipt ref",
    "✓ Validate every QR type (MEMBER · ENT · TKT · DRINK · BDAY · generic)",
    "✓ Explain base tier vs subscription tier",
    "✓ Walk a member to their Member Card",
    "✓ Triage a missing-points complaint",
    "✓ Explain (don't fire) the SOS workflow",
    "✓ Identify when to escalate to manager / Luna Hub",
], 0.7, 2.2, 12, 5, size=17)
footer_band(s, "Module 9 · Sign-off")

# ---------- Closing slide ----------
s = add_blank_slide(); fill_bg(s, DEEP)
add_text(s, "WELCOME TO THE TEAM", 0.7, 2.6, 12, 1.2,
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Now go make every Luna night unforgettable.", 0.7, 4.2, 12, 0.8,
         size=20, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "Questions?  ·  support@lunagroupapp.com.au", 0.7, 6.3, 12, 0.4,
         size=12, color=MUTED, align=PP_ALIGN.CENTER)

prs.save(PPTX_OUT)
print(f"  → {PPTX_OUT} ({os.path.getsize(PPTX_OUT)//1024} KB)")
print(f"  Slides: {len(prs.slides)}")

print("\n✓ Done")
